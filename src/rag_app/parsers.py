import csv
import io
import os
import tempfile
from dataclasses import dataclass
from pathlib import Path

import PyPDF2
import docx
from bs4 import BeautifulSoup
from pptx import Presentation

from .docling_utils import convert_bytes_with_docling, pdf_needs_full_ocr

USE_DOCLING = os.getenv("RAG_USE_DOCLING", "1") != "0"


@dataclass
class ParsedDocument:
    text: str
    parser: str = "plain"
    used_ocr: bool = False


def parse_text_file(content: bytes) -> str:
    return content.decode("utf-8", errors="ignore")


def parse_pdf_pypdf(content: bytes) -> str:
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(content))
        pages = []
        for page in reader.pages:
            try:
                pages.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n\n".join(pages)
    except Exception:
        return ""


def parse_pdf_docling(content: bytes, force_full_page_ocr: bool = False) -> str:
    return convert_bytes_with_docling(content, ".pdf", force_full_page_ocr=force_full_page_ocr)


def parse_pdf(content: bytes) -> ParsedDocument:
    force_ocr = pdf_needs_full_ocr(content)
    if USE_DOCLING:
        try:
            text = parse_pdf_docling(content, force_full_page_ocr=force_ocr)
            if text.strip():
                parser = "docling-ocr" if force_ocr else "docling"
                return ParsedDocument(text=text, parser=parser, used_ocr=force_ocr)
        except Exception:
            pass
    return ParsedDocument(text=parse_pdf_pypdf(content), parser="pypdf2", used_ocr=False)


def parse_docx(content: bytes) -> str:
    document = docx.Document(io.BytesIO(content))
    paragraphs = [para.text for para in document.paragraphs]
    return "\n\n".join(paragraphs)


def parse_docx_docling(content: bytes) -> ParsedDocument:
    if not USE_DOCLING:
        return ParsedDocument(text=parse_docx(content), parser="python-docx")
    try:
        text = convert_bytes_with_docling(content, ".docx", force_full_page_ocr=False)
        if text.strip():
            return ParsedDocument(text=text, parser="docling")
    except Exception:
        pass
    return ParsedDocument(text=parse_docx(content), parser="python-docx")


def parse_pptx(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    slides = []
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append("\n".join([text for text in texts if text.strip()]))
    return "\n\n".join(slides)


def parse_html(content: bytes) -> str:
    decoded = content.decode("utf-8", errors="ignore")
    soup = BeautifulSoup(decoded, "html.parser")
    return soup.get_text(separator="\n\n")


def parse_csv(content: bytes) -> str:
    text = content.decode("utf-8", errors="ignore")
    rows = []
    reader = csv.reader(text.splitlines())
    for row in reader:
        rows.append(", ".join(row))
    return "\n".join(rows)


def parse_upload_bytes(filename: str, data: bytes) -> ParsedDocument:
    name = filename.lower()
    if name.endswith(".pdf"):
        return parse_pdf(data)
    if name.endswith(".docx"):
        return parse_docx_docling(data)
    if name.endswith(".pptx"):
        return ParsedDocument(text=parse_pptx(data), parser="python-pptx")
    if name.endswith(".html") or name.endswith(".htm"):
        return ParsedDocument(text=parse_html(data), parser="html")
    if name.endswith(".txt") or name.endswith(".md"):
        return ParsedDocument(text=parse_text_file(data), parser="plain")
    if name.endswith(".csv"):
        return ParsedDocument(text=parse_csv(data), parser="csv")
    return ParsedDocument(text=parse_text_file(data), parser="plain")


def parse_upload_file(uploaded) -> str:
    """Backward-compatible API — returns plain text."""
    return parse_upload_bytes(uploaded.name, uploaded.getvalue()).text
